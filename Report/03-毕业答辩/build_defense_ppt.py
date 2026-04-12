from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "Report/02-中期报告/鄢玺中期报告.pptx"
OUTPUT = ROOT / "Report/03-毕业答辩/鄢玺毕业答辩.pptx"

SPEEDUP = ROOT / "Img/mybenchmark/final-comparison-speedup.png"
HEATMAP = ROOT / "Img/mybenchmark/final-comparison-heatmap.png"
STRIPPLOT = ROOT / "Img/mybenchmark/statistical-validation-stripplot.png"

EMU_PER_INCH = 914400

DARK = "003366"
NAVY = "002A57"
GREEN = "DCEFD8"
TEAL = "DDF6F6"
BG = "F7FBFB"
WHITE = "FFFFFF"
GRAY = "5B6470"
RED = "C0504D"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def emu_to_inches(value: int) -> float:
    return value / EMU_PER_INCH


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def remove_all_slides_except_first(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)[1:]
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def set_text_frame(
    text_frame,
    paragraphs: list[dict],
    *,
    margin: float = 0.08,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = valign
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)

    for idx, item in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_before = Pt(item.get("space_before", 0))
        paragraph.space_after = Pt(item.get("space_after", 5))
        if "line_spacing" in item:
            paragraph.line_spacing = item["line_spacing"]
        if "level" in item:
            paragraph.level = item["level"]

        runs = item.get("runs")
        if runs is None:
            runs = [
                {
                    "text": item.get("text", ""),
                    "size": item.get("size", 16),
                    "bold": item.get("bold", False),
                    "color": item.get("color", DARK),
                }
            ]

        for run_spec in runs:
            run = paragraph.add_run()
            run.text = run_spec["text"]
            run.font.size = Pt(run_spec.get("size", item.get("size", 16)))
            run.font.bold = run_spec.get("bold", item.get("bold", False))
            run.font.italic = run_spec.get("italic", item.get("italic", False))
            run.font.underline = run_spec.get("underline", item.get("underline", False))
            run.font.color.rgb = rgb(run_spec.get("color", item.get("color", DARK)))
            if "name" in run_spec:
                run.font.name = run_spec["name"]


def set_shape_text(
    shape,
    paragraphs: list[dict],
    *,
    margin: float = 0.08,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    set_text_frame(shape.text_frame, paragraphs, margin=margin, valign=valign)


def set_placeholder_text(
    slide,
    idx: int,
    paragraphs: list[dict],
    *,
    margin: float = 0.08,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    shape = slide.placeholders[idx]
    set_shape_text(shape, paragraphs, margin=margin, valign=valign)
    return shape


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    paragraphs: list[dict],
    *,
    fill: str | None = None,
    line: str | None = None,
    radius: bool = False,
    margin: float = 0.08,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    if fill is None and line is None and not radius:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            if radius
            else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        )
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill or WHITE)
        shape.line.color.rgb = rgb(line or fill or WHITE)
        shape.line.width = Pt(1.2)

    set_shape_text(shape, paragraphs, margin=margin, valign=valign)
    return shape


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    with Image.open(path) as image:
        aspect = image.width / image.height

    box_ratio = w / h
    if aspect >= box_ratio:
        pic_w = w
        pic_h = w / aspect
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * aspect
        pic_x = x + (w - pic_w) / 2
        pic_y = y

    slide.shapes.add_picture(
        str(path),
        Inches(pic_x),
        Inches(pic_y),
        width=Inches(pic_w),
        height=Inches(pic_h),
    )


def add_picture_in_placeholder(slide, idx: int, path: Path) -> None:
    placeholder = slide.placeholders[idx]
    add_picture_contain(
        slide,
        path,
        emu_to_inches(placeholder.left),
        emu_to_inches(placeholder.top),
        emu_to_inches(placeholder.width),
        emu_to_inches(placeholder.height),
    )


def add_slide_title(slide, title: str, *, size: int = 28) -> None:
    set_shape_text(
        slide.shapes.title,
        [{"text": title, "size": size, "bold": True, "color": DARK}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_section_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_shape_text(
        slide.shapes.title,
        [
            {
                "text": title,
                "size": 30,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.LEFT,
            }
        ],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        1,
        [{"text": subtitle, "size": 18, "color": GRAY}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return slide


def add_assertion_band(
    slide,
    text: str,
    *,
    y: float = 1.12,
    fill: str = NAVY,
    text_color: str = WHITE,
) -> None:
    add_textbox(
        slide,
        0.82,
        y,
        9.55,
        0.42,
        [
            {
                "text": text,
                "size": 13,
                "bold": True,
                "color": text_color,
                "align": PP_ALIGN.CENTER,
            }
        ],
        fill=fill,
        line=fill,
        radius=True,
        margin=0.02,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    value: str,
    *,
    fill: str,
    value_color: str = NAVY,
) -> None:
    add_textbox(
        slide,
        x,
        y,
        w,
        h,
        [
            {
                "text": label,
                "size": 13,
                "bold": True,
                "color": DARK,
                "align": PP_ALIGN.CENTER,
                "space_after": 2,
            },
            {
                "text": value,
                "size": 20,
                "bold": True,
                "color": value_color,
                "align": PP_ALIGN.CENTER,
            },
        ],
        fill=fill,
        line=fill,
        radius=True,
        margin=0.05,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_flow_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: str,
) -> None:
    add_textbox(
        slide,
        x,
        y,
        w,
        h,
        [
            {
                "text": title,
                "size": 17,
                "bold": True,
                "color": DARK,
                "align": PP_ALIGN.CENTER,
                "space_after": 6,
            },
            {
                "text": body,
                "size": 13,
                "color": DARK,
                "align": PP_ALIGN.CENTER,
                "line_spacing": 1.1,
            },
        ],
        fill=fill,
        line=fill,
        radius=True,
        margin=0.09,
        valign=MSO_ANCHOR.MIDDLE,
    )


def set_table_cell(
    cell,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: str = DARK,
    align: PP_ALIGN = PP_ALIGN.CENTER,
) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.03)
    text_frame.margin_bottom = Inches(0.03)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def bullet(text: str, *, size: int = 16, bold: bool = False, color: str = DARK) -> dict:
    return {"text": text, "size": size, "bold": bold, "color": color}


def build_deck() -> None:
    prs = Presentation(str(TEMPLATE))
    remove_all_slides_except_first(prs)

    prs.core_properties.author = "Yanxi"
    prs.core_properties.title = "基于LoongArch架构的ART执行引擎适配与优化技术研究"
    prs.core_properties.subject = "硕士毕业论文答辩"

    # 1. 封面：直接回归中期第一页的样式路径
    slide = prs.slides[0]
    slide.shapes.title.left = Inches(1.15)
    slide.shapes.title.top = Inches(1.78)
    slide.shapes.title.width = Inches(9.0)
    slide.shapes.title.height = Inches(1.55)
    set_shape_text(
        slide.shapes.title,
        [
            {
                "text": "基于LoongArch架构的ART执行引擎\n适配与优化技术研究",
                "size": 34,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
                "line_spacing": 1.0,
            }
        ],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    slide.placeholders[1].text = ""
    add_textbox(
        slide,
        1.65,
        4.72,
        8.1,
        1.55,
        [
            {
                "text": "鄢玺",
                "size": 30,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
                "space_after": 5,
            },
            {
                "text": "导师：张福新",
                "size": 30,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
                "space_after": 5,
            },
            {
                "text": "2026年6月",
                "size": 32,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
            },
        ],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 2. 提纲
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_slide_title(slide, "答辩提纲")
    set_placeholder_text(
        slide,
        1,
        [
            {
                "runs": [
                    {"text": "01  ", "size": 19, "bold": True, "color": NAVY},
                    {"text": "研究背景与目标", "size": 19, "bold": True, "color": DARK},
                ],
                "space_after": 12,
            },
            {
                "runs": [
                    {"text": "02  ", "size": 19, "bold": True, "color": NAVY},
                    {"text": "执行引擎适配与验证", "size": 19, "bold": True, "color": DARK},
                ],
                "space_after": 12,
            },
            {
                "runs": [
                    {"text": "03  ", "size": 19, "bold": True, "color": NAVY},
                    {"text": "性能优化研究与结果", "size": 19, "bold": True, "color": DARK},
                ],
                "space_after": 12,
            },
            {
                "runs": [
                    {"text": "04  ", "size": 19, "bold": True, "color": NAVY},
                    {"text": "总结与展望", "size": 19, "bold": True, "color": DARK},
                ]
            },
        ],
        margin=0.18,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 3. 分节页
    add_section_slide(prs, "一、研究背景与目标", "问题定义、研究主线与论文贡献")

    # 4. 问题定义与主要贡献
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "问题定义与主要贡献")
    set_placeholder_text(
        slide,
        1,
        [{"text": "核心问题", "size": 19, "bold": True, "color": NAVY}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        2,
        [
            bullet("AOSP 15 官方主线尚未提供 LoongArch 平台 ART 执行引擎实现。"),
            bullet("论文必须同时回答“能运行、能验证、能优化”三个问题。"),
            bullet("性能结论需要统一 baseline 与可复核证据链支撑。"),
        ],
        margin=0.10,
    )
    set_placeholder_text(
        slide,
        3,
        [{"text": "论文回答", "size": 19, "bold": True, "color": NAVY}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        4,
        [
            bullet("完成执行引擎关键模块适配，建立 LoongArch 平台 ART 执行闭环。"),
            bullet("建立编译、部署、run-test、libcore、benchmark 的统一验证链路。"),
            bullet("筛选三项正式方案，并区分稳定收益、边界性收益与失败样本。"),
        ],
        margin=0.10,
    )
    add_textbox(
        slide,
        0.9,
        6.15,
        9.55,
        0.42,
        [
            {
                "text": "主线不是零散 patch，而是“适配设计 → 验证收敛 → 优化筛选”的系统闭环。",
                "size": 13,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.CENTER,
            }
        ],
        fill=BG,
        line=GREEN,
        radius=True,
        margin=0.03,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 5. 研究路线
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    add_slide_title(slide, "整体研究路线")
    set_placeholder_text(
        slide,
        1,
        [
            {
                "text": "本文从“可运行”推进到“可验证”再到“可优化”，最终给出稳定主结论与收益边界。",
                "size": 16,
                "bold": True,
                "color": NAVY,
                "space_after": 8,
            },
            bullet("执行引擎适配解决平台可用性问题。"),
            bullet("统一验证链路保证语义正确与性能比较口径一致。"),
            bullet("正式单次对比与重复测量共同约束优化结论的强度。"),
        ],
        margin=0.10,
    )
    flow_placeholder = slide.placeholders[2]
    fx = emu_to_inches(flow_placeholder.left)
    fy = emu_to_inches(flow_placeholder.top)
    fw = emu_to_inches(flow_placeholder.width)
    card_w = 2.15
    card_h = 1.18
    gap = (fw - card_w * 4) / 3
    card_y = fy - 0.04
    flow_specs = [
        ("执行引擎适配", "补齐关键模块，建立 LoongArch 平台可运行 ART。", GREEN),
        ("统一验证链路", "以 run-test / libcore / benchmark 建立可复核证据。", TEAL),
        ("热点分析筛选", "固定 baseline，定位真正影响 workload 的主导热点。", GREEN),
        ("正式结论形成", "区分稳定收益、边界性收益与失败样本。", TEAL),
    ]
    for idx, (title, body, fill) in enumerate(flow_specs):
        x = fx + idx * (card_w + gap)
        add_flow_card(slide, x, card_y, card_w, card_h, title, body, fill=fill)
        if idx < len(flow_specs) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(x + card_w + 0.08),
                Inches(card_y + 0.40),
                Inches(gap - 0.16),
                Inches(0.38),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(NAVY)
            arrow.line.color.rgb = rgb(NAVY)

    add_metric_card(slide, 1.05, 5.72, 2.8, 0.68, "工程产出", "执行闭环", fill=GREEN)
    add_metric_card(slide, 4.05, 5.72, 2.8, 0.68, "方法产出", "统一验证链路", fill=TEAL)
    add_metric_card(slide, 7.05, 5.72, 2.8, 0.68, "研究产出", "稳定主收益结论", fill=GREEN)

    # 6. 分节页
    add_section_slide(prs, "二、执行引擎适配与验证", "建立可运行、可验证的 LoongArch 平台 ART")

    # 7. 适配闭环
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide, "执行引擎适配的系统闭环")
    add_assertion_band(
        slide,
        "核心目标：以执行路径为主线补齐关键模块，形成可运行、可验证的 LoongArch 平台 ART 执行闭环",
        y=1.18,
    )
    add_textbox(
        slide,
        0.95,
        1.82,
        9.45,
        0.55,
        [
            {
                "text": "从汇编器、解释器到 JIT、JNI 与运行时入口，关键不是单点可跑，而是整条执行链条同时闭合。",
                "size": 15,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.CENTER,
                "line_spacing": 1.05,
            }
        ],
        fill=BG,
        line=WHITE,
        radius=True,
        margin=0.02,
        valign=MSO_ANCHOR.MIDDLE,
    )
    module_specs = [
        (
            0.95,
            2.72,
            "Assembler 与基础生成",
            "寄存器、指令编码、分支与链接修补能力。",
            GREEN,
        ),
        (
            4.00,
            2.72,
            "C++ 解释器与 Nterp",
            "补齐解释执行与字节码派发的关键路径。",
            TEAL,
        ),
        (
            7.05,
            2.72,
            "JIT 后端与 OSR",
            "打通 HIR 到 LoongArch64 机器码生成。",
            GREEN,
        ),
        (
            0.95,
            4.32,
            "JNI 编译器与桥接",
            "闭合 Java / native 调用链与栈帧约定。",
            TEAL,
        ),
        (
            4.00,
            4.32,
            "Runtime Fast Path",
            "补齐字符串、调用桩、deopt 等运行时入口。",
            GREEN,
        ),
        (
            7.05,
            4.32,
            "Intrinsic 与能力扩展",
            "补充标准库热点路径与执行引擎成熟度。",
            TEAL,
        ),
    ]
    for x, y, title, body, fill in module_specs:
        add_flow_card(slide, x, y, 2.45, 1.18, title, body, fill=fill)
    add_textbox(
        slide,
        0.95,
        6.00,
        9.45,
        0.42,
        [
            {
                "text": "结论：适配结果是执行引擎系统闭环，不是若干体系结构相关文件的零散移植。",
                "size": 13,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.CENTER,
            }
        ],
        fill=BG,
        line=GREEN,
        radius=True,
        margin=0.02,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 8. 系统验证结果
    slide = prs.slides.add_slide(prs.slide_layouts[14])
    add_slide_title(slide, "正确性验证结果")
    table_shape = slide.placeholders[1].insert_table(7, 3)
    table = table_shape.table
    table.columns[0].width = Inches(2.35)
    table.columns[1].width = Inches(1.65)
    table.columns[2].width = Inches(6.25)

    headers = ["验证项", "结果", "说明"]
    rows = [
        ("run-test 用例总数", "1033", "art/test 当前 run-test 用例总数"),
        (
            "compiler 变体",
            "7 类",
            "interp-ac、interpreter、jit、jit-on-first-use、optimizing、speed-profile、baseline",
        ),
        ("理论配置数", "7231", "1033 × 7"),
        ("自动 skip", "430", "根据 knownfailures 规则自动跳过"),
        ("实际执行数", "6801", "系统级目标机回归实际执行集合"),
        ("libcore", "主要测试包通过", "少量因架构条件限制自动跳过，不影响主要正确性结论"),
    ]
    for row in table.rows:
        row.height = Inches(0.65)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(NAVY)
        set_table_cell(cell, header, size=13, bold=True, color=WHITE)

    for row_idx, row in enumerate(rows, start=1):
        fill_color = BG if row_idx % 2 else WHITE
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill_color)
            align = PP_ALIGN.LEFT if col_idx == 2 else PP_ALIGN.CENTER
            size = 11.5 if col_idx == 2 else 12
            set_table_cell(cell, value, size=size, bold=False, color=DARK, align=align)

    # 9. 分节页
    add_section_slide(prs, "三、性能优化研究", "统一 baseline、热点驱动与收益边界")

    # 10. 单次结果总览
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    add_slide_title(slide, "三项正式方案的单次结果总览")
    add_textbox(
        slide,
        0.84,
        2.02,
        4.55,
        4.38,
        [
            {
                "text": "统一 baseline 下，同日三项正式方案的收益层级已经足够清晰。",
                "size": 15,
                "bold": True,
                "color": NAVY,
                "space_after": 8,
            },
            bullet("方案一：JIT 热度阈值调整，pmd +3.30%，但复测稳定性不足。"),
            bullet("方案二：字符串搬运批量化，pmd +0.86%，同时 lu.small -8.89%。"),
            bullet("方案三：LSX 浮点 SIMD，lu.small +44.40%，且证据链最完整。"),
        ],
        margin=0.10,
    )
    add_picture_contain(slide, SPEEDUP, 5.48, 2.02, 5.55, 4.42)

    # 11. 稳定主结论
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    add_slide_title(slide, "稳定主结论：LSX 浮点 SIMD 向量化")
    add_textbox(
        slide,
        0.84,
        2.02,
        4.55,
        4.20,
        [
            {
                "text": "LSX 浮点 SIMD 是当前唯一同时满足收益、解释性与复现性的正式方案。",
                "size": 15,
                "bold": True,
                "color": NAVY,
                "space_after": 8,
            },
            bullet("补齐 float32/float64 的核心 lowering，直接命中 lu.small 主热点。"),
            bullet("正式单次结果：46.22 → 66.74 ops/m，增益 +44.40%。"),
            bullet("同日 3 轮复测均值 +45.70%，标准差仅 0.50 ops/m。"),
        ],
        margin=0.10,
    )
    add_metric_card(slide, 0.95, 5.56, 1.92, 0.76, "单次增益", "+44.40%", fill=GREEN)
    add_metric_card(slide, 3.05, 5.56, 2.10, 0.76, "复测均值", "+45.70%", fill=TEAL)
    add_picture_contain(slide, STRIPPLOT, 5.48, 2.02, 5.55, 4.42)

    # 12. 收益边界
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    add_slide_title(slide, "收益边界与方法论结论")
    set_placeholder_text(
        slide,
        1,
        [
            {
                "text": "优化是否值得保留，取决于收益边界是否清晰、证据链是否完整。",
                "size": 15,
                "bold": True,
                "color": NAVY,
                "space_after": 8,
            },
            bullet("方案一在 pmd 上存在性能提升样本，但复测均值尚未稳定转正。"),
            bullet("方案二说明“局部热点下降”不等于“全局 benchmark 提升”。"),
            bullet("方案三收益集中在规则浮点循环，并不会自动推广到所有数值 workload。"),
            bullet("因此，LoongArch 平台 ART 优化应优先命中主导热点，并严格控制副作用。"),
        ],
        margin=0.10,
    )
    add_picture_in_placeholder(slide, 2, HEATMAP)

    # 13. 分节页
    add_section_slide(prs, "四、总结与展望", "论文结论、方法价值与后续工作")

    # 14. 总结与展望
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "总结与展望")
    set_placeholder_text(
        slide,
        1,
        [{"text": "论文主结论", "size": 19, "bold": True, "color": NAVY}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        2,
        [
            bullet("AOSP 15 ART 在 LoongArch 平台上的执行引擎适配是可行的。"),
            bullet("统一验证链路与固定 baseline 是性能研究具备解释价值的前提。"),
            bullet("只有命中主导热点、边界清晰且可复现的方案才值得保留。"),
        ],
        margin=0.10,
    )
    set_placeholder_text(
        slide,
        3,
        [{"text": "后续工作", "size": 19, "bold": True, "color": NAVY}],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        4,
        [
            bullet("完善 AOT / 安装期编译支持。"),
            bullet("扩展 LSX / LASX 能力与更高层门控。"),
            bullet("引入更多 benchmark 与应用级验证。"),
            bullet("增强低扰动观测与运行时因果分析。"),
        ],
        margin=0.10,
    )
    add_textbox(
        slide,
        0.92,
        6.12,
        9.55,
        0.38,
        [
            {
                "text": "本文的核心产出不仅是工程闭环，更是“什么值得保留、为什么成立、边界在哪里”的研究方法。",
                "size": 13,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.CENTER,
            }
        ],
        fill=BG,
        line=GREEN,
        radius=True,
        margin=0.02,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 15. 致谢
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_shape_text(
        slide.shapes.title,
        [
            {
                "text": "感谢聆听！",
                "size": 30,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.LEFT,
            }
        ],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_placeholder_text(
        slide,
        1,
        [
            {
                "text": "请各位老师批评指正",
                "size": 20,
                "bold": True,
                "color": DARK,
                "space_after": 8,
            },
            {"text": "Q&A", "size": 16, "color": GRAY},
        ],
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
